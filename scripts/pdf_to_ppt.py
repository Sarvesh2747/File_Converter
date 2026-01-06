
import sys
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
import os

def convert_pdf_to_ppt(pdf_path, ppt_path):
    print(f"Converting {pdf_path} -> {ppt_path}")
    
    # Create presentation
    prs = Presentation()
    # Blank slide layout
    BLANK_SLIDE_LAYOUT = 6
    
    doc = fitz.open(pdf_path)
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        
        # Render page to image
        pix = page.get_pixmap(dpi=150)
        img_filename = f"temp_page_{page_num}.png"
        pix.save(img_filename)
        
        # Add slide
        slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
        
        # Calculate aspect ratio
        # Default slide is 10x7.5 inches
        # We can adjust slide size to match PDF page?
        # For now, let's fit the image to the slide
        
        left = top = Inches(0)
        # Using A4 ratio approx or just fit
        # Ideally we should resize presentation to match PDF, but let's simple fit
        slide.shapes.add_picture(img_filename, left, top, height=prs.slide_height)
        
        # Cleanup temp image
        os.remove(img_filename)
        
    prs.save(ppt_path)
    print("SUCCESS")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python pdf_to_ppt.py <input_pdf> <output_ppt>")
        sys.exit(1)
        
    input_file = sys.argv[1]
    output_file = sys.argv[2]
    
    try:
        convert_pdf_to_ppt(input_file, output_file)
    except Exception as e:
        print(f"Error: {e}")
        sys.exit(1)
